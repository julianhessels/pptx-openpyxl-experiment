from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

pr1 = Presentation()

slide1_layout = pr1.slide_layouts[0]

slide1 = pr1.slides.add_slide(slide1_layout)

title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]

title1.text= "ANALYSTRISING"
subtitle1.text = "Subscribe to my channel"

slide2_layout = pr1.slide_layouts[1]
slide2 = pr1.slides.add_slide(slide2_layout)

title2 = slide2.shapes.title
title2.text = "Now For Some Bullet Points"

bullet_point_box = slide2.shapes
bullet_points_lvl1 = bullet_point_box.placeholders[1]
bullet_points_lvl1.text ="Subscribe"

shapes = slide1.shapes
print(shapes)



bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl2.text ="to"
bullet_points_lvl2.level = 1

bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl3.text ="my"
bullet_points_lvl3.level = 2

bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl3.text ="Channel!"
bullet_points_lvl3.level = 3

#Add Slide 3
slide3_layout = pr1.slide_layouts[5]
slide3 = pr1.slides.add_slide(slide3_layout)
title3 = slide3.shapes.title
title3.text = "Picture Time!"

img1 = "Elements.jpg"
from_left = Inches(3)
from_top = Inches(4)
add_picture = slide3.shapes.add_picture(img1,from_left,from_top)

#Part3
#Add Slide 4
slide4_layout = pr1.slide_layouts[5]
slide4 = pr1.slides.add_slide(slide4_layout)
title4 = slide4.shapes.title
title4.text = "Shapework"

#Create a shape
left1 = top1 = width1 = height1 = Inches(2)
add_shape1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left1,top1,width1,height1)

left2 = Inches(6)
top2 = Inches(2)
width2 = height2 = Inches(2)
arrow1 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,left2,top2,width2,height2)


fill_arrow1 = arrow1.fill
fill_arrow1.solid()
fill_arrow1.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_5
arrow1.text = "Pijl111"

arrow1.rotation = 90

#pr1.save('GreatPresentation.pptx')
#pr1.save('GreatPresentation_Part2.pptx')
pr1.save('GreatPresentation_Part3.pptx')


